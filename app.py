import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import io
import base64
import tempfile
import os
import json
import requests
from PIL import Image
import random
import re

st.set_page_config(page_title="AI PowerPoint Maker", layout="wide")

# App title and description
st.title("AI PowerPoint Maker")
st.subheader("Create beautiful presentations with AI")

# Load API key from secrets
api_key = None
try:
    api_key = st.secrets["GEMINI_API_KEY"]
except Exception:
    st.error("API Key not found in secrets. Please set up your API key in .streamlit/secrets.toml file.")
    st.info("Add the following line to your .streamlit/secrets.toml file: GEMINI_API_KEY = 'your-api-key-here'")

# Sidebar for configuration
with st.sidebar:
    st.header("Configuration")
    
    st.subheader("Presentation Settings")
    num_slides = st.slider("Number of slides", min_value=3, max_value=15, value=5)
    
    presentation_style = st.selectbox(
        "Presentation Style",
        ["Professional", "Creative", "Minimalist", "Academic", "Business", "Tech Pitch", "Professional"]
    )
    
    color_theme = st.selectbox(
        "Color Theme",
        ["Blue", "Green", "Red", "Purple", "Orange", "Grayscale", "Teal Accent", "Midnight", "Grayscale Pro", 
         "VISME Modern", "VISME Gradient", "VISME Corporate"]
    )
    
    st.subheader("Visual Settings")
    # Enhanced animation controls like VISME
    animation_style = st.selectbox(
        "Animation Style",
        ["None", "Fade", "Wipe", "Fly In", "Split", "Zoom", "Pop In", "Zoom & Fade", "Sleek Entrance", 
         "VISME Sequential", "VISME Cascading", "VISME Staggered", "VISME Dynamic"]
    )
    
    # Animation speed control
    animation_speed = st.select_slider(
        "Animation Speed",
        options=["Very Slow", "Slow", "Medium", "Fast", "Very Fast"],
        value="Medium"
    )
    
    include_images = st.checkbox("Include AI-generated image placeholders", value=True)
    
    slide_transition = st.selectbox(
        "Slide Transition",
        ["None", "Fade", "Push", "Reveal", "Morph", "Subtle Zoom", 
         "VISME Smooth", "VISME Elegant", "VISME Dynamic"]
    )
    
    modern_elements = st.checkbox("Add modern design elements", value=True)
    
    # New design enhancement options
    design_complexity = st.slider("Design Complexity", min_value=1, max_value=5, value=3,
                               help="Higher values add more decorative elements")
    
    use_gradients = st.checkbox("Use gradient backgrounds", value=True, 
                             help="Adds professional gradient backgrounds to slides")
    
    shape_style = st.selectbox(
        "Shape Style",
        ["Rounded", "Angular", "Organic", "Geometric", "Minimal"]
    )

# Main content area
st.header("Create Your Presentation")
topic = st.text_area("Enter your presentation topic or describe what you need", height=100)
additional_info = st.text_area("Additional instructions or specific points to include (optional)", height=100)

# Improved text color and visibility function
def ensure_text_readability(text_frame, theme_color, is_title=False, base_size=18, title_size=32):
    """
    Dynamically adjust text color for better readability against backgrounds
    
    Args:
        text_frame: PowerPoint text frame
        theme_color: Color dictionary for the current theme
        is_title: Boolean to differentiate title from body text
        base_size: Default text size for body text
        title_size: Default text size for titles
    """
    # Calculate background luminance
    bg_color = theme_color["subtle"]
    luminance = 0.299 * bg_color[0] + 0.587 * bg_color[1] + 0.114 * bg_color[2]
    
    # Choose text color based on background luminance
    if luminance > 128:
        # Light background: use dark text
        text_color = theme_color["text"]
        if is_title:
            text_color = theme_color["main"]
    else:
        # Dark background: use light text
        text_color = (255, 255, 255)
        if is_title:
            text_color = theme_color["accent"]
    
    # Apply the calculated text color and size
    for paragraph in text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(*text_color)
        paragraph.font.size = Pt(title_size if is_title else base_size)
        if is_title:
            paragraph.font.bold = True

# Function to generate presentation content using Gemini
def generate_presentation_content(topic, additional_info, num_slides, style):
    if not api_key:
        st.error("API key is missing. Please configure it in the secrets.toml file.")
        return None
    
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-2.0-flash')
        
        prompt = f"""
        Create a professional, visually striking presentation on: {topic}
        Additional information: {additional_info}
        Style: {style}
        Number of slides: {num_slides}
        
        Design this presentation using modern Canva/Visme principles:
        - Strong visual hierarchy with varying text sizes
        - Concise, impactful bullet points (max 5-6 per slide)
        - Strategic use of white space
        - Engaging headlines (6-8 words maximum)
        - Visual storytelling flow from slide to slide
        
        Format the response as a JSON with this structure:
        {{
            "title": "Main presentation title (short, impactful, max 5-7 words)",
            "slides": [
                {{
                    "title": "Slide title (clear, concise, engaging)",
                    "content": "Slide content in brief bullet points (25-40 words max per bullet)",
                    "type": "title_slide/content_slide/image_slide/split_slide/quote_slide/statistic_slide",
                    "image_description": "Detailed description for generating an image related to this slide",
                    "layout_note": "Brief design suggestion for this particular slide (optional)"
                }},
                ...
            ]
        }}
        
        For best results:
        - Vary slide types for visual interest (title_slide, content_slide, image_slide, split_slide, quote_slide, statistic_slide)
        - Make title slides bold and minimal (large text, minimal bullets)
        - Keep content slides focused on 1 key message per slide
        - For statistic slides, highlight just 1-2 key numbers with minimal text
        - Include thoughtful transition slides between major sections
        - End with a clear, actionable conclusion slide
        
        The presentation should tell a coherent visual story from beginning to end, not just present information.
        """
        
        response = model.generate_content(prompt)
        
        # Parse the response to get JSON
        response_text = response.text
        # Extract JSON part if needed
        if "```json" in response_text:
            json_str = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            json_str = response_text.split("```")[1].strip()
        else:
            json_str = response_text
            
        # Clean and parse JSON
        presentation_data = json.loads(json_str)
        return presentation_data
        
    except Exception as e:
        st.error(f"Error generating presentation content: {str(e)}")
        return None

def create_powerpoint(presentation_data, style, color_theme, animation_style="None", slide_transition="None"):
    prs = Presentation()
    
    # Modern font selections
    TITLE_FONT = "Montserrat"
    BODY_FONT = "Open Sans"
    
    # Enhanced color palette with modern gradients and more sophisticated colors
    color_themes = {
        "Professional Blue": {
            "gradient_start": (41, 128, 185),     # Vibrant blue
            "gradient_end": (109, 213, 250),      # Soft light blue
            "text_primary": (33, 47, 61),         # Dark blue-gray
            "text_secondary": (52, 152, 219),     # Bright blue accent
            "background": (240, 248, 255)         # Very light blue-white
        },
        "Modern Minimalist": {
            "gradient_start": (55, 59, 68),       # Dark gray
            "gradient_end": (66, 134, 244),       # Bright blue
            "text_primary": (28, 28, 30),         # Almost black
            "text_secondary": (108, 117, 125),    # Muted gray
            "background": (248, 249, 250)         # Very light gray
        },
        "Tech Gradient": {
            "gradient_start": (25, 32, 50),       # Deep navy
            "gradient_end": (106, 130, 251),      # Bright periwinkle
            "text_primary": (255, 255, 255),      # Pure white
            "text_secondary": (173, 181, 189),    # Light gray
            "background": (33, 37, 41)            # Near-black
        }
    }
    
    # Select theme or default to Professional Blue
    theme = color_themes.get(color_theme, color_themes["Professional Blue"])
    
    # Utility function for creating gradient background
    def create_gradient_background(slide):
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor(*theme["gradient_start"])
        fill.gradient_stops[1].color.rgb = RGBColor(*theme["gradient_end"])
        fill.gradient_angle = 45  # Diagonal gradient
    
    # Enhanced text styling function
    def style_text(text_frame, is_title=False, is_subtitle=False):
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = TITLE_FONT if is_title else BODY_FONT
            paragraph.font.bold = is_title
            
            if is_title:
                paragraph.font.size = Pt(36)  # Large, impactful title
                paragraph.font.color.rgb = RGBColor(*theme["text_primary"])
            elif is_subtitle:
                paragraph.font.size = Pt(20)  # Slightly smaller subtitle
                paragraph.font.color.rgb = RGBColor(*theme["text_secondary"])
            else:
                paragraph.font.size = Pt(16)  # Clean, readable body text
                paragraph.font.color.rgb = RGBColor(*theme["text_primary"])
            
            # Center alignment for titles, left for body
            paragraph.alignment = PP_ALIGN.CENTER if is_title else PP_ALIGN.LEFT
    
    # Ensure presentation data is not None
    if not presentation_data:
        st.error("No presentation data available.")
        return None
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    create_gradient_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Safely get title, use a default if None
    title_text = presentation_data.get("title", "AI-Generated Presentation")
    title.text = title_text
    subtitle.text = "AI-Powered Presentation"
    
    style_text(title.text_frame, is_title=True)
    style_text(subtitle.text_frame, is_subtitle=True)
    
    # Content slides generation
    slides_data = presentation_data.get('slides', [])
    
    # Ensure we have at least one slide
    if not slides_data:
        # Create a default content slide if no slides are provided
        slides_data = [{
            'title': 'Default Slide',
            'content': 'No content was generated. Please try again.',
            'type': 'content_slide'
        }]
    
    for slide_data in slides_data:
        # Choose appropriate layout
        slide_layout = {
            'title_slide': prs.slide_layouts[0],
            'content_slide': prs.slide_layouts[1],
            'section_slide': prs.slide_layouts[2],
            'two_content_slide': prs.slide_layouts[3]
        }.get(slide_data.get('type', 'content_slide'), prs.slide_layouts[1])
        
        slide = prs.slides.add_slide(slide_layout)
        create_gradient_background(slide)
        
        # Slide title
        if slide.shapes.title:
            # Safely get slide title, use a default if None
            slide_title = slide_data.get('title', 'Untitled Slide')
            slide.shapes.title.text = slide_title
            style_text(slide.shapes.title.text_frame, is_title=True)
        
        # Slide content
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            
            # Safely handle content, convert to string, use default if None
            slide_content = str(slide_data.get('content', 'No content available'))
            
            # Trim content if it's extremely long
            if len(slide_content) > 500:
                slide_content = slide_content[:500] + '...'
            
            content_placeholder.text = slide_content
            style_text(content_placeholder.text_frame)
        
        # Optional: Add modern design elements
        if slide.shapes.title:
            # Add a subtle accent line under the title
            x, y = slide.shapes.title.left, slide.shapes.title.top + slide.shapes.title.height + Inches(0.1)
            width = slide.shapes.title.width
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                x, y, 
                width, Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*theme["text_secondary"])
    
    # Temporary file path for saving
    temp_pptx_path = os.path.join(tempfile.gettempdir(), 'modern_presentation.pptx')
    prs.save(temp_pptx_path)
    
    return temp_pptx_path

# Function to convert the presentation to base64 for download
def get_binary_file_downloader_html(bin_file_path, file_label='File'):
    with open(bin_file_path, 'rb') as f:
        data = f.read()
    b64 = base64.b64encode(data).decode()
    href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{file_label}.pptx">Download {file_label}</a>'
    return href

# Create button to generate presentation
if st.button("Create PowerPoint Presentation"):
    if not topic:
        st.error("Please enter a presentation topic.")
    else:
        with st.spinner("Generating presentation content with AI..."):
            presentation_data = generate_presentation_content(
                topic, 
                additional_info, 
                num_slides, 
                presentation_style
            )
            
        if presentation_data:
            with st.spinner("Creating PowerPoint file..."):
                pptx_path = create_powerpoint(
                    presentation_data, 
                    presentation_style, 
                    color_theme, 
                    animation_style,
                    slide_transition
                )
                
                if pptx_path:
                    st.success("Presentation created successfully!")
                    
                    # Display download link
                    file_name = "AI_Generated_Presentation"
                    st.markdown(get_binary_file_downloader_html(pptx_path, file_name), unsafe_allow_html=True)
                    
                    # Summary of presentation
                    st.subheader("Presentation Summary")
                    st.write(f"**Title:** {presentation_data['title']}")
                    st.write(f"**Slides:** {len(presentation_data['slides'])} slides created")
                    st.write(f"**Style:** {presentation_style} with {color_theme} color theme")
                    
                    # Preview first slide as text
                    st.subheader("Preview of First Slide")
                    if presentation_data['slides']:
                        first_slide = presentation_data['slides'][0]
                        st.write(f"**Title:** {first_slide['title']}")
                        st.write(f"**Content:** {first_slide['content'][:200]}...")
                    
                    # Clean up temp file after a delay
                    import threading
                    def cleanup():
                        import time
                        time.sleep(300)  # Wait 5 minutes before cleaning up
                        if os.path.exists(pptx_path):
                            os.unlink(pptx_path)
                    
                    threading.Thread(target=cleanup).start()
        else:
            st.error("Failed to generate presentation content. Please try again.")

# Add footer
st.markdown("---")
st.markdown("AI PowerPoint Maker | Made with Streamlit and Gemini AI")